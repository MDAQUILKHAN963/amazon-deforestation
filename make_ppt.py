"""Generate a clean, professional 5-slide PowerPoint summarizing the project."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette --------------------------------------------------------------
GREEN      = RGBColor(0x1B, 0x5E, 0x20)   # primary
GREEN_LT   = RGBColor(0x43, 0xA0, 0x47)   # accent
MINT       = RGBColor(0xE8, 0xF5, 0xE9)   # light fill
DARK       = RGBColor(0x21, 0x21, 0x21)
GREY       = RGBColor(0x60, 0x60, 0x60)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
FONT = "Calibri"


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=False):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = color
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    """runs: list of (text, size, bold, color)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (t, sz, b, c) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        r = p.add_run(); r.text = t
        r.font.name = FONT; r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c
    return box


def bullets(slide, items, x=1.0, y=2.0, w=11.3, h=4.6, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        t, lvl = (item if isinstance(item, tuple) else (item, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(10)
        if not t:
            continue
        bullet = "●  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = bullet; r.font.name = FONT
        r.font.size = Pt(size - 2 * lvl); r.font.bold = True
        r.font.color.rgb = GREEN_LT if lvl == 0 else GREY
        r2 = p.add_run(); r2.text = t; r2.font.name = FONT
        r2.font.size = Pt(size - 2 * lvl); r2.font.color.rgb = DARK if lvl == 0 else GREY
    return box


def header(slide, title, idx):
    rect(slide, 0, 0, SW.inches, 1.25, GREEN)                 # top band
    rect(slide, 0, 1.25, SW.inches, 0.07, GREEN_LT)          # accent line
    text(slide, 0.7, 0, 11.5, 1.25, [(title, 30, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    # footer
    rect(slide, 0, 7.18, SW.inches, 0.32, MINT)
    text(slide, 0.7, 7.16, 9, 0.32, [("Amazon Deforestation Detection", 11, False, GREEN)],
         anchor=MSO_ANCHOR.MIDDLE)
    text(slide, 11.6, 7.16, 1.1, 0.32, [(f"{idx} / 7", 11, True, GREEN)],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ===== Slide 1 — Title =====================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW.inches, SH.inches, GREEN)                    # full green bg
rect(s, 0, 3.95, SW.inches, 0.06, GREEN_LT)
text(s, 1.0, 1.85, 11.3, 2.0, [
    ("Amazon Deforestation Detection", 44, True, WHITE),
    ("Detecting deforested areas from satellite imagery using deep learning", 21, False, MINT),
], space=12)
text(s, 1.0, 4.25, 11.3, 0.6, [("Internship Project — Deep Learning / Computer Vision", 17, True, WHITE)])
# presenter / details block
text(s, 1.0, 5.05, 11.3, 2.0, [
    ("Presented by:  Mohammad Aquil Khan", 18, True, WHITE),
    ("Jaypee University of Engineering and Technology", 16, False, MINT),
    ("Mentor:  Dr. Abdul Qadir", 16, False, MINT),
    ("Internship:  [add internship name]", 16, False, MINT),
], space=8)

# ===== Slide 2 — Problem & Goal ===========================================
s = prs.slides.add_slide(BLANK)
header(s, "Problem & Goal", 2)
bullets(s, [
    "Amazon deforestation is a major environmental problem that must be monitored.",
    "Manually checking satellite images over huge areas is slow and impractical.",
    "Goal: automatically detect deforested areas in satellite imagery.",
    "Task type: image segmentation — classify every pixel as forest or deforested.",
    ("Focused on a study region in the Amazon for the years 2019–2021.", 1),
], y=2.0)

# ===== Slide 3 — Dataset ==================================================
s = prs.slides.add_slide(BLANK)
header(s, "Dataset", 3)
bullets(s, [
    "Used the public MultiEarth 2023 Amazon satellite dataset.",
    "Sentinel-2 — optical (color) imagery; bands B2, B3, B4, B8.",
    "Sentinel-1 — radar imagery (VV, VH) that can see through clouds.",
    "Labels: hand-drawn masks marking deforested areas (the ground truth).",
    "Prepared ~4,700 image tiles of 256×256 pixels, each with 6 channels.",
    ("Matched each label to satellite imagery within a ±2 month window.", 1),
    ("Removed cloudy images and normalized pixel values.", 1),
], y=1.8, size=19)

# ===== Slide 4 — Approach / Methodology ===================================
s = prs.slides.add_slide(BLANK)
header(s, "Approach", 4)
bullets(s, [
    "Model: U-Net with an EfficientNet backbone (encoder–decoder for segmentation).",
    "Encoder understands the image; decoder rebuilds a pixel-by-pixel deforestation map.",
    "Input = 6 channels (4 optical + 2 radar) so the model sees beyond color alone.",
    "Output = 1 map: each pixel labeled forest or deforested.",
    "Location-based train/validation split so the model is tested on unseen areas.",
    ("Data augmentation (flips, rotations) for more variety and better generalization.", 1),
], y=1.9)

# ===== Slide 5 — How the Model Learns =====================================
s = prs.slides.add_slide(BLANK)
header(s, "How the Model Learns (Training)", 5)
bullets(s, [
    "Training repeats a simple 4-step cycle on small batches of images:",
    ("Forward — the model predicts a deforestation map.", 1),
    ("Loss — compare the prediction to the true mask; get an error score.", 1),
    ("Backward — find what caused the error (backpropagation).", 1),
    ("Update — the optimizer adjusts the model to reduce the error.", 1),
    "Trained for 10 epochs (full passes) on a single free GPU.",
    "Kept the best version of the model based on validation score.",
], y=1.8, size=19)

# ===== Slide 6 — Demo (Web App) ==========================================
s = prs.slides.add_slide(BLANK)
header(s, "Demo — Web App", 6)
bullets(s, [
    "Built a simple web app so anyone can test the model without writing code.",
    "User picks a location and year → the model analyzes that satellite tile.",
    "Shows three images side by side:",
    ("Satellite image (what the area looks like).", 1),
    ("Actual deforestation (the real labeled map).", 1),
    ("Model's prediction (what the model detected).", 1),
    "Gives a clear verdict and the percentage of area deforested.",
], y=1.9, size=19)

# ===== Slide 7 — Progress & Next Steps ===================================
s = prs.slides.add_slide(BLANK)
header(s, "Progress & Next Steps", 7)
text(s, 1.0, 1.7, 11, 0.5, [("Completed so far", 22, True, GREEN)])
bullets(s, [
    ("Collected and prepared the satellite dataset (~4,700 tiles).", 1),
    ("Built the U-Net + EfficientNet model.", 1),
    ("Set up and ran the full training pipeline.", 1),
    ("Built a working web app to demo predictions.", 1),
], y=2.3, size=19)
text(s, 1.0, 4.2, 11, 0.5, [("Next steps", 22, True, GREEN)])
bullets(s, [
    ("Fine-tune and improve the model further.", 1),
    ("Support more locations using live satellite data (Google Earth Engine).", 1),
    ("Extend to forest fire / burned-area detection.", 1),
], y=4.8, size=19)

out = "Amazon_Deforestation_Project_progress.pptx"
prs.save(out)
print("saved", out)
